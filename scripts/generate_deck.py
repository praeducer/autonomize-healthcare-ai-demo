"""Generate the Autonomize AI PA Solution Architecture PPTX deck.

SSOT chain:
  solution-architecture.md  (ground truth)
    -> presentation.md       (narrative SSOT)
      -> this file            (slide-optimized for PPTX)
        -> PPTX               (rendered output)

When editing: update presentation.md first, then sync changes here.

Usage:   python scripts/generate_deck.py
Shortcut: make deck
"""

from __future__ import annotations

import struct
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
PROJECT_ROOT = Path(__file__).resolve().parent.parent
DIAGRAMS = PROJECT_ROOT / "docs" / "architecture" / "diagrams"
OUTPUT_DIR = PROJECT_ROOT / "docs" / "presentation"
LOGO_WHITE = PROJECT_ROOT / "docs" / "assets" / "autonomize-logo.png"  # for dark backgrounds
LOGO_DARK = PROJECT_ROOT / "docs" / "assets" / "autonomize-logo-dark.png"  # for light backgrounds

# ---------------------------------------------------------------------------
# Brand theme — extracted from autonomize.ai CSS
# ---------------------------------------------------------------------------
PRIMARY = RGBColor(0x29, 0x06, 0x5A)  # deep purple
ACCENT = RGBColor(0x73, 0x1E, 0xE3)  # bright purple
LIGHT_BG = RGBColor(0xF8, 0xF4, 0xFE)  # pale lavender
TEXT_COLOR = RGBColor(0x1A, 0x1A, 0x1A)  # near-black
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT_NAME = "Calibri"

# ---------------------------------------------------------------------------
# Layout constants (widescreen 13.33" x 7.5")
# ---------------------------------------------------------------------------
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
MARGIN = Inches(1.0)
CONTENT_W = Inches(11.33)
TITLE_H = Inches(0.9)
CONTENT_Y = Inches(1.05)
FOOTER_H = Inches(0.35)
GAP = Inches(0.12)
IMG_WITH_TABLE = 2.8  # inches — for slides with image + table below

# Font sizes
FONT_TITLE = Pt(26)
FONT_HEADING = Pt(15)
FONT_BODY = Pt(12)
FONT_CALLOUT = Pt(11)
FONT_TABLE = Pt(10)


def _png_dimensions(path: Path) -> tuple[int, int]:
    """Read PNG width/height from file header (bytes 16-23)."""
    with open(path, "rb") as f:
        f.seek(16)
        w, h = struct.unpack(">II", f.read(8))
    return w, h


def _add_title_bar(slide, title: str) -> None:
    """Colored rectangle with white title text + white logo."""
    shape = slide.shapes.add_shape(
        1,
        Inches(0),
        Inches(0),
        SLIDE_W,
        TITLE_H,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = FONT_TITLE
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = FONT_NAME
    p.alignment = PP_ALIGN.LEFT

    # White logo in top-right of title bar
    if LOGO_WHITE.exists():
        logo_h = Inches(0.35)
        pw, ph = _png_dimensions(LOGO_WHITE)
        logo_w = Inches(0.35 * pw / ph)
        slide.shapes.add_picture(
            str(LOGO_WHITE),
            Emu(SLIDE_W.emu - MARGIN.emu - logo_w.emu),
            Emu(int((TITLE_H.emu - logo_h.emu) / 2)),
            logo_w,
            logo_h,
        )


def _add_footer(slide, num: int, total: int) -> None:
    """Slide number + dark logo in footer."""
    tx_box = slide.shapes.add_textbox(
        SLIDE_W - MARGIN - Inches(1.5),
        SLIDE_H - FOOTER_H,
        Inches(1.5),
        FOOTER_H,
    )
    tf = tx_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num} / {total}"
    p.font.size = Pt(8)
    p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    p.font.name = FONT_NAME
    p.alignment = PP_ALIGN.RIGHT

    # Dark logo in bottom-left
    if LOGO_DARK.exists():
        logo_h = Inches(0.2)
        pw, ph = _png_dimensions(LOGO_DARK)
        logo_w = Inches(0.2 * pw / ph)
        slide.shapes.add_picture(
            str(LOGO_DARK),
            MARGIN,
            Emu(SLIDE_H.emu - FOOTER_H.emu + Inches(0.05).emu),
            logo_w,
            logo_h,
        )
    p.alignment = PP_ALIGN.RIGHT


def _remaining(y_emu: int) -> int:
    """Available vertical space from y to footer."""
    return SLIDE_H.emu - FOOTER_H.emu - y_emu


def _estimate_lines(text: str, width_inches: float, font_pt: float) -> int:
    """Estimate how many lines text will occupy at a given font size and width."""
    # ~9 chars/inch at 12pt Calibri proportional, scales inversely with font size
    chars_per_inch = 72 / font_pt * 1.3
    chars_per_line = max(1, int(width_inches * chars_per_inch))
    # Count explicit newlines + wrapping
    total = 0
    for segment in text.split("\n"):
        total += max(1, -(-len(segment) // chars_per_line))  # ceil division
    return total


def _add_text_block(slide, text: str, y: int, font_size=FONT_BODY) -> int:
    """Add a text box. Returns new y position (EMU)."""
    if not text.strip():
        return y + Inches(0.1).emu  # minimal spacer for empty text
    line_h = int(font_size.emu * 1.5)  # line height ~150% of font size
    num_lines = _estimate_lines(text, CONTENT_W.inches, font_size.pt)
    h = max(Inches(0.2).emu, num_lines * line_h)
    h = min(h, _remaining(y))
    if h < Inches(0.15).emu:
        return y

    tx_box = slide.shapes.add_textbox(MARGIN, Emu(y), CONTENT_W, Emu(h))
    tf = tx_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = TEXT_COLOR
    p.font.name = FONT_NAME
    return y + h + GAP.emu


def _add_heading_block(slide, text: str, y: int, font_size=FONT_HEADING) -> int:
    """Add a heading text box. Returns new y."""
    h = Inches(0.35).emu
    if h > _remaining(y):
        return y

    tx_box = slide.shapes.add_textbox(MARGIN, Emu(y), CONTENT_W, Emu(h))
    tf = tx_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.font.name = FONT_NAME
    return y + h + GAP.emu


def _add_image_block(slide, path: str, y: int, max_h: float | None = None) -> int:
    """Add a centered image. Returns new y."""
    img_path = Path(path)
    pw, ph = _png_dimensions(img_path)
    aspect = ph / pw

    img_w = CONTENT_W.inches
    img_h = img_w * aspect

    # Cap height
    content_max = _remaining(y) / 914400  # convert EMU to inches
    cap = min(content_max, max_h or content_max * 0.55)
    if img_h > cap:
        img_h = cap
        img_w = img_h / aspect

    x_offset = (CONTENT_W.inches - img_w) / 2
    slide.shapes.add_picture(
        str(img_path),
        Emu(MARGIN.emu + int(x_offset * 914400)),
        Emu(y),
        Inches(img_w),
        Inches(img_h),
    )
    return y + int(img_h * 914400) + GAP.emu


def _add_callout_block(slide, text: str, y: int, label: str = "", font_size=FONT_CALLOUT) -> int:
    """Add a callout box with accent bar. Returns new y."""
    full_text = (f"{label}: " if label else "") + text
    line_h = int(font_size.emu * 1.5)
    num_lines = _estimate_lines(full_text, CONTENT_W.inches - 0.4, font_size.pt)
    padding = Inches(0.1).emu
    h = max(Inches(0.3).emu, num_lines * line_h + padding * 2)
    h = min(h, _remaining(y))
    if h < Inches(0.2).emu:
        return y

    # Background rectangle
    bg = slide.shapes.add_shape(1, MARGIN, Emu(y), CONTENT_W, Emu(h))
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_BG
    bg.line.fill.background()

    # Accent bar
    bar = slide.shapes.add_shape(1, MARGIN, Emu(y), Inches(0.06), Emu(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    # Text — vertically centered in the callout box
    tx_box = slide.shapes.add_textbox(
        Emu(MARGIN.emu + Inches(0.15).emu),
        Emu(y),
        Emu(CONTENT_W.emu - Inches(0.3).emu),
        Emu(h),
    )
    tf = tx_box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    # Vertically center text in the callout box (OOXML bodyPr anchor attribute)
    tf._txBody.bodyPr.set(qn("a:anchor"), "ctr")
    p = tf.paragraphs[0]
    if label:
        run_label = p.add_run()
        run_label.text = f"{label}: "
        run_label.font.bold = True
        run_label.font.size = font_size
        run_label.font.color.rgb = PRIMARY
        run_label.font.name = FONT_NAME
        run_body = p.add_run()
        run_body.text = text
        run_body.font.size = font_size
        run_body.font.color.rgb = TEXT_COLOR
        run_body.font.name = FONT_NAME
    else:
        p.text = text
        p.font.size = font_size
        p.font.color.rgb = TEXT_COLOR
        p.font.name = FONT_NAME

    return y + h + GAP.emu


def _add_table_block(
    slide,
    headers: list[str],
    rows: list[list[str]],
    y: int,
    font_size: int = 10,
) -> int:
    """Add a table. Returns new y."""
    num_cols = len(headers)
    avail = _remaining(y)

    # Adaptive row height
    total_rows = len(rows) + 1
    ideal_row_h = Inches(0.35).emu
    min_row_h = Inches(0.25).emu
    row_h = max(min_row_h, min(ideal_row_h, avail // max(total_rows, 1)))
    max_visible = max(1, avail // row_h - 1)
    visible_rows = rows[:max_visible]
    num_rows = len(visible_rows) + 1  # +1 for header

    table_h = num_rows * row_h
    if table_h > avail:
        table_h = avail

    shape = slide.shapes.add_table(num_rows, num_cols, MARGIN, Emu(y), CONTENT_W, Emu(table_h))
    tbl = shape.table

    # Column widths — proportional to content length
    max_lens = []
    for ci in range(num_cols):
        col_max = len(headers[ci])
        for row in visible_rows:
            if ci < len(row):
                col_max = max(col_max, len(row[ci]))
        max_lens.append(max(col_max, 1))
    total_len = sum(max_lens)

    for ci in range(num_cols):
        proportion = max_lens[ci] / total_len
        tbl.columns[ci].width = Emu(int(proportion * CONTENT_W.emu))

    # Header row
    for ci, h_text in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.text = h_text
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(font_size)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = FONT_NAME
            p.alignment = PP_ALIGN.LEFT
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Data rows
    for ri, row_data in enumerate(visible_rows):
        for ci in range(num_cols):
            cell = tbl.cell(ri + 1, ci)
            cell.text = row_data[ci] if ci < len(row_data) else ""
            # Alternating row shading
            if ri % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.color.rgb = TEXT_COLOR
                p.font.name = FONT_NAME
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    return y + table_h + GAP.emu


# ---------------------------------------------------------------------------
# Slide data — maps 1:1 to presentation.md
# ---------------------------------------------------------------------------
def _diagram(name: str) -> str:
    return str(DIAGRAMS / name)


SLIDES: list[dict] = [
    {
        "title": "AI-Driven Prior Authorization",
        "layout": "title",  # special layout: vertically centered
        "body": [
            {"type": "heading", "text": "Solution Architecture for a Large US Health Plan", "font_size": 20},
            {"type": "text", "text": "Paul Prae  |  Principal AI Engineer & Architect", "font_size": 16},
            {"type": "text", "text": "www.paulprae.com", "font_size": 14},
        ],
    },
    {
        "title": "The Problem & Opportunity",
        "body": [
            {"type": "heading", "text": "The Problem"},
            {
                "type": "text",
                "text": "Manual PA processing costs $10.97 per provider transaction and $3.52 per payer transaction, versus $0.05 when fully electronic (CAQH 2024). 93% of physicians report PA delays patient care (AMA 2024).",
            },
            {"type": "heading", "text": "The Opportunity \u2014 Altais + Autonomize AI (Feb 2026)"},
            {
                "type": "table",
                "headers": ["Metric", "Result"],
                "rows": [
                    ["PA review time reduction", "45%"],
                    ["Manual error reduction", "54%"],
                    ["Auto-determination rate", "50%"],
                ],
            },
            {"type": "heading", "text": "Why Autonomize AI"},
            {
                "type": "text",
                "text": "\u2022 Live at 3 of the 5 largest US health plans\n\u2022 PA Copilot on Azure Marketplace \u2014 Microsoft Pegasus Program member\n\u2022 ServiceNow partnership extends reach into payer workflows\n\u2022 CMS-0057-F deadline (Jan 2027) creates urgency",
            },
        ],
    },
    {
        "title": "PA Request Lifecycle",
        "body": [
            {"type": "image", "path": _diagram("03-pa-request-flow.png"), "max_h": 4.0},
            {
                "type": "callout",
                "text": "Submit \u2192 Intake (OCR) \u2192 Validate (eligibility) \u2192 AI Review (coverage matching + confidence) \u2192 Route (auto-approve / human review / pend) \u2192 Respond (payer core + provider notification)",
            },
        ],
    },
    {
        "title": "Demo \u2014 Proof of Concept",
        "body": [
            {
                "type": "callout",
                "label": "LIVE DEMO",
                "text": "Working proof of concept \u2014 demonstrates core clinical review (steps 4-5 of the lifecycle).",
            },
            {"type": "image", "path": _diagram("07-demo-architecture.png"), "max_h": IMG_WITH_TABLE},
            {
                "type": "table",
                "font_size": 9,
                "headers": ["Demo Scope (Phase 0)", "Production Differences"],
                "rows": [
                    ["5 PA cases with real ICD-10/CPT codes", "Mock eligibility \u2192 Payer Core integration"],
                    ["AI reviews in ~30s via Claude tool use", "Local criteria \u2192 CMS Coverage Database API"],
                    [
                        "FHIR R4 data models (fhir.resources R4B)",
                        "Synthetic data \u2192 real clinical records via FHIR",
                    ],
                    ["CLI, REST API + Swagger, Web Dashboard", "Single-user \u2192 Azure Container Apps + auto-scale"],
                ],
            },
        ],
    },
    {
        "title": "System Context",
        "body": [
            {"type": "image", "path": _diagram("01-system-context.png"), "max_h": 3.2},
            {
                "type": "table",
                "headers": ["Actor", "Role", "Integration"],
                "rows": [
                    ["Healthcare Providers", "Submit PA requests", "Fax, Portal, X12 278"],
                    ["Autonomize AI Platform", "AI-driven clinical review", "PA Copilot on Genesis"],
                    ["Health Plan Systems", "Eligibility, benefits, clinical data", "REST API, FHIR R4"],
                    ["Regulators (CMS)", "Compliance reporting", "CMS-0057-F metrics"],
                ],
            },
        ],
    },
    {
        "title": "Solution Architecture",
        "body": [
            {"type": "image", "path": _diagram("02-component-architecture.png"), "max_h": 5.5},
        ],
    },
    {
        "title": "Solution Architecture \u2014 Components",
        "body": [
            {
                "type": "table",
                "font_size": 10,
                "headers": ["Component", "Azure Service", "Purpose"],
                "rows": [
                    ["Ingestion Gateway", "API Mgmt + Functions", "Receives all PA channels"],
                    ["Document Processing", "AI Document Intelligence", "OCR for faxes"],
                    ["Eligibility Service", "Payer Core REST API", "Member validation"],
                    ["Clinical Data Aggregation", "Health Data Services (FHIR R4)", "Unified clinical context"],
                    ["PA Copilot", "Genesis Platform + Claude", "AI clinical review"],
                    ["Determination Router", "Functions + Rules", "Confidence-based routing"],
                    ["Clinical Review Dashboard", "Autonomize AI Studio", "Human reviewer interface"],
                    ["Determination Response", "Azure Functions", "Payer Core writeback"],
                    ["LLMOps Pipeline", "Azure Monitor + Evals", "Performance monitoring"],
                    ["Audit & Compliance", "Immutable Blob Storage", "Tamper-proof audit trail"],
                ],
            },
        ],
    },
    {
        "title": "Why This Architecture",
        "body": [
            {
                "type": "text",
                "font_size": 13,
                "text": (
                    "\u2022 Safety-first routing \u2014 confidence thresholds route low-certainty cases to human reviewers; no auto-deny without clinical review\n\n"
                    "\u2022 Configurable thresholds \u2014 start conservative, tune with real performance data\n\n"
                    "\u2022 CMS-0057-F compliance \u2014 FHIR R4 foundation meets Jan 2027 API deadline without re-architecture\n\n"
                    "\u2022 Azure-native \u2014 leverages Autonomize\u2019s Azure ecosystem, Pegasus Program, and Marketplace presence\n\n"
                    "\u2022 Full audit trail \u2014 every determination records model version, input hash, reasoning, evidence, and confidence for 7-year retention"
                ),
            },
        ],
    },
    {
        "title": "Top 3 Security Risks & Mitigations",
        "body": [
            {
                "type": "table",
                "font_size": 9,
                "headers": ["#", "Risk", "Architectural Mitigation", "Operational Mitigation"],
                "rows": [
                    [
                        "1",
                        "PHI exposure through AI pipeline",
                        "PHI tokenization before LLM \u2014 AI sees clinical facts without patient identity",
                        "Zero data retention for model training (enterprise terms)",
                    ],
                    [
                        "2",
                        "Prompt injection via clinical docs",
                        "Document sanitization + system prompt isolation",
                        "Output validation requires evidence citations",
                    ],
                    [
                        "3",
                        "Untraceable AI decisions",
                        "Tamper-proof audit: model version, input hash, reasoning, confidence",
                        "Immutable 7-year retention, append-only writes",
                    ],
                ],
            },
            {
                "type": "callout",
                "label": "Additional controls",
                "text": "Entra ID RBAC, AES-256 at rest, TLS 1.2+ in transit, private endpoints, no auto-deny without human review.",
            },
        ],
    },
    {
        "title": "Progressive Delivery",
        "body": [
            {
                "type": "table",
                "headers": ["Phase", "Focus", "Key Deliverable"],
                "rows": [
                    ["Phase 0: Demo", "Prove the concept", "Working AI PA review with mock data"],
                    ["Phase 1: MVP", "Single LOB, single channel", "Production PA processing with human review"],
                    ["Phase 2: Scale", "Multi-channel, multi-LOB", "Fax OCR, legacy data, LOB configuration"],
                    ["Phase 3: Enterprise", "Full scale, compliance", "All channels, 20 LOBs, CMS reporting"],
                ],
            },
            {
                "type": "callout",
                "label": "Design principle",
                "text": "Each phase produces a deployable, demonstrable system. Decision gates between phases use real performance data to scope the next phase.",
            },
        ],
    },
    {
        "title": "Discussion Starters",
        "body": [
            {"type": "heading", "text": "Business Strategy"},
            {
                "type": "text",
                "text": "\u2022 How does the ServiceNow partnership change payer integration strategy?\n\u2022 What is the target auto-determination rate for Phase 1?",
            },
            {"type": "heading", "text": "Technical Depth"},
            {
                "type": "text",
                "text": "\u2022 How does the Genesis Platform handle coverage criteria updates?\n\u2022 What\u2019s the Azure AI Foundry Agent Service integration status?",
            },
            {"type": "heading", "text": "Implementation"},
            {
                "type": "text",
                "text": "\u2022 Which LOB is the ideal Phase 1 candidate?\n\u2022 What\u2019s been the biggest integration challenge with existing payer deployments?",
            },
        ],
    },
    {
        "title": "Appendix A: Clinical Data Integration",
        "body": [
            {"type": "image", "path": _diagram("04-clinical-data-access.png"), "max_h": IMG_WITH_TABLE},
            {
                "type": "table",
                "headers": ["Source", "Protocol", "Auth"],
                "rows": [
                    ["Modern EMRs", "FHIR R4 REST API", "OAuth 2.0 / SMART on FHIR"],
                    ["Legacy DB Connector", "DB connector / HL7 v2", "Service account + VNet"],
                ],
            },
            {
                "type": "callout",
                "label": "FHIR R4",
                "text": "Interoperability standard for clinical data. Modern sources expose natively; legacy data normalized before AI processing.",
            },
            {
                "type": "callout",
                "label": "Security",
                "text": "All clinical data passes through PHI tokenization before reaching the AI engine.",
            },
        ],
    },
    {
        "title": "Appendix B: AI Model Monitoring & Feedback",
        "body": [
            {"type": "image", "path": _diagram("06-llmops-pipeline.png"), "max_h": IMG_WITH_TABLE},
            {"type": "heading", "text": "Detect drift"},
            {
                "type": "text",
                "text": "\u2022 Overturn rate monitoring (human overrides AI), appeal rate, accuracy trends\n\u2022 Automated evals \u2014 golden test cases benchmarked on schedule\n\u2022 Confidence distribution shifts signal model or data changes",
            },
            {"type": "heading", "text": "Feedback loop"},
            {
                "type": "text",
                "text": "1. Reviewer corrections \u2192 updated eval dataset  2. Benchmark new vs current  3. Staged blue-green rollout  4. Guardrails always active",
            },
        ],
    },
    {
        "title": "Appendix C: Scaling to 20 LOBs",
        "body": [
            {
                "type": "table",
                "headers": ["Approach", "Cost", "Isolation", "Complexity"],
                "rows": [
                    ["Multi-tenant", "Lower", "Logical", "Lower"],
                    ["Multi-instance", "Higher", "Physical", "Higher"],
                ],
            },
            {
                "type": "callout",
                "label": "Recommendation",
                "text": "Start multi-tenant with per-LOB config. Genesis Platform supports it. Separate instances only where regulation requires physical isolation.",
            },
            {
                "type": "callout",
                "label": "Honest unknowns",
                "text": "The right answer depends on actual LOB rule complexity and regulatory requirements \u2014 both are discovery questions.",
            },
        ],
    },
    {
        "title": "Sources & References",
        "body": [
            {
                "type": "table",
                "font_size": 8,
                "headers": ["#", "Claim", "Source"],
                "rows": [
                    ["1", "$10.97 manual PA labor cost (provider)", "CAQH 2024 Index, p. 9"],
                    ["2", "$3.52 manual PA labor cost (payer)", "CAQH 2024 Index, p. 9"],
                    ["3", "~$0.05 electronic PA cost (payer)", "CAQH 2024 Index, p. 9"],
                    ["4", "45% PA review time reduction", "Altais + Autonomize AI (BusinessWire Feb 2026)"],
                    ["5", "54% manual error reduction", "Altais + Autonomize AI (BusinessWire Feb 2026)"],
                    ["6", "50% auto-determination rate", "Altais + Autonomize AI (BusinessWire Feb 2026)"],
                    ["7", "93% physicians say PA delays care", "AMA 2024 Prior Auth Survey, p. 5"],
                    ["8", "CMS-0057-F Phase 1 Jan 2026", "CMS Interoperability Final Rule Fact Sheet"],
                    ["9", "CMS-0057-F Phase 2 Jan 2027", "CMS Interoperability Final Rule Fact Sheet"],
                    ["10", "Live at 3 of 5 largest US plans", "Autonomize AI (BusinessWire Feb 2026)"],
                    ["11", "PA Copilot on Azure Marketplace", "Azure Marketplace listing"],
                    ["12", "ServiceNow partnership", "BusinessWire Mar 2026"],
                    ["13", "Claude in Azure AI Foundry", "Microsoft Learn"],
                    ["14", "Genesis Platform", "GlobeNewsWire Apr 2025"],
                    ["15", "Microsoft Pegasus Program", "BusinessWire Nov 2025"],
                ],
            },
        ],
    },
]


# ---------------------------------------------------------------------------
# Render
# ---------------------------------------------------------------------------
def render_pptx(slides: list[dict], output_path: Path) -> None:
    """Generate a PPTX file from structured slide data."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Use a blank layout
    blank_layout = prs.slide_layouts[6]  # blank

    total = len(slides)
    for si, slide_data in enumerate(slides):
        slide = prs.slides.add_slide(blank_layout)
        _add_title_bar(slide, slide_data["title"])
        _add_footer(slide, si + 1, total)

        # Title slide: vertically center the content block
        is_title_layout = slide_data.get("layout") == "title"
        if is_title_layout:
            y = int(SLIDE_H.emu * 0.35)  # start at ~35% down for visual center
        else:
            y = CONTENT_Y.emu

        for block in slide_data["body"]:
            if _remaining(y) < Inches(0.15).emu:
                break

            btype = block["type"]
            if btype == "text":
                fs = Pt(block.get("font_size", FONT_BODY.pt))
                y = _add_text_block(slide, block["text"], y, font_size=fs)
            elif btype == "heading":
                fs = Pt(block.get("font_size", FONT_HEADING.pt))
                y = _add_heading_block(slide, block["text"], y, font_size=fs)
            elif btype == "image":
                y = _add_image_block(slide, block["path"], y, max_h=block.get("max_h"))
            elif btype == "callout":
                y = _add_callout_block(
                    slide,
                    block["text"],
                    y,
                    label=block.get("label", ""),
                    font_size=Pt(block.get("font_size", FONT_CALLOUT.pt)),
                )
            elif btype == "table":
                y = _add_table_block(
                    slide,
                    block["headers"],
                    block["rows"],
                    y,
                    font_size=block.get("font_size", FONT_TABLE.pt),
                )

    prs.save(str(output_path))
    print(f"PPTX -> {output_path}")


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def main() -> None:
    now = datetime.now()
    stamp = now.strftime("%Y-%m-%d-%H%M")
    base = f"autonomize-ai-pa-solution-architecture-paul-prae-{stamp}"

    print(f"\nGenerating {len(SLIDES)} slides -> {base}\n")
    render_pptx(SLIDES, OUTPUT_DIR / f"{base}.pptx")
    print("\nDone.\n")


if __name__ == "__main__":
    main()
